"""Generate the 10-slide executive PowerPoint."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt

SLIDES=[("Healthcare Patient Dataset","Machine Learning, Visualization & Responsible Model Governance"),("2. Business Problem","Classify Test Results while deciding whether synthetic data supports responsible deployment."),("3. Dataset","55,500 rows • 15 columns • 0 missing • 534 duplicates • 108 negative billing records"),("4. Exploratory Analysis","Uniform categories, weak correlations, duplicates, negative billing, and text-quality issues."),("5. Modeling","Baseline Logistic Regression versus advanced CatBoost; identifiers excluded to limit memorization."),("6. Results","Logistic Regression accuracy 0.323; CatBoost accuracy 0.317; chance level ≈ 0.333."),("7. Key Insights","Additional model complexity did not create meaningful predictive signal."),("8. Recommendations","Use for portfolio learning; do not use for clinical or patient-level decisions."),("9. Future Work","Acquire governed real-world data; add temporal validation, calibration, fairness, and monitoring."),("10. Conclusion","Responsible data science includes recognizing when a model should not be deployed.")]

def main():
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
    for title,body in SLIDES:
        slide=prs.slides.add_slide(prs.slide_layouts[6])
        p=slide.shapes.add_textbox(Inches(.8),Inches(.8),Inches(11.8),Inches(1.1)).text_frame.paragraphs[0]; p.text=title; p.font.size=Pt(30); p.font.bold=True
        p2=slide.shapes.add_textbox(Inches(.8),Inches(2.1),Inches(11.6),Inches(3.5)).text_frame.paragraphs[0]; p2.text=body; p2.font.size=Pt(22)
    out=Path(__file__).with_name("Healthcare_Patient_Dataset_Executive_Presentation.pptx"); prs.save(out); print(f"Saved {out}")

if __name__=="__main__": main()
